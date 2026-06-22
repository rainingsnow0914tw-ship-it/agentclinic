"""Build the AgentClinic Devpost submission deck from the UiPath
hackathon template (docs/Submission deck.pptx).

Output: docs/Submission deck - AgentClinic.pptx

Strategy: open the template, walk each slide's shapes by name, replace
placeholder text in-place (preserves font/color/size formatting from
the template's first run). For slide 2 (team), keep one team member
block and delete the other three.
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
SRC = ROOT / "docs" / "Submission deck.pptx"
DST = ROOT / "docs" / "Submission deck - AgentClinic.pptx"


def set_text(text_frame, new_text: str) -> None:
    """Replace a single-line text frame, preserving first-run formatting."""
    p = text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = new_text
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.add_run().text = new_text
    for extra in list(text_frame.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)


def set_multiline(text_frame, lines: list[str]) -> None:
    """Replace a text frame with multiple paragraphs (preserves first-run
    formatting on first paragraph, copies style for the rest)."""
    p0 = text_frame.paragraphs[0]
    # remove extra paragraphs first
    for extra in list(text_frame.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    # set first line
    if p0.runs:
        p0.runs[0].text = lines[0]
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.add_run().text = lines[0]
    # add the rest
    for line in lines[1:]:
        new_p = text_frame.add_paragraph()
        new_p.text = line


def remove_shape(shape) -> None:
    sp = shape._element
    sp.getparent().remove(sp)


prs = Presentation(SRC)

# ============================================================
# SLIDE 1 · Cover
# ============================================================
s1 = prs.slides[0]
for shape in s1.shapes:
    if not shape.has_text_frame:
        continue
    txt = shape.text_frame.text
    if "Presentation title" in txt:
        # Title placeholder — replace with our title + subtitle on 2 lines
        set_multiline(shape.text_frame, [
            "AgentClinic",
            "Forensic analysis for AI agents, native to UiPath Test Cloud.",
        ])
    # Keep the "UiPath AgentHack. Build the AI agents of tomorrow." line

# ============================================================
# SLIDE 2 · Team (Individual — keep one, delete the other three)
# ============================================================
s2 = prs.slides[1]
# Inventory: 4 Picture frames (174/177/180/183) + 4 Name shapes
# (175/178/181/184) + 4 Title shapes (176/179/182/185) + 1 Team/Project
# shape (186). Keep 174 (first picture), 175 (first name), 176 (first
# title), 186 (team+project). Delete the rest.
KEEP_IDS = {174, 175, 176, 186}
shapes_to_remove = []
for shape in s2.shapes:
    if shape.shape_id not in KEEP_IDS:
        shapes_to_remove.append(shape)
for shape in shapes_to_remove:
    remove_shape(shape)

# Now fill the kept shapes
for shape in s2.shapes:
    if not shape.has_text_frame:
        continue
    if shape.shape_id == 175:  # Name
        set_text(shape.text_frame, "Chloe Kao")
    elif shape.shape_id == 176:  # Title + email
        set_multiline(shape.text_frame, [
            "Independent Builder",
            "rainingsnow0914.tw@gmail.com",
        ])
    elif shape.shape_id == 186:  # Team name + Project title
        set_multiline(shape.text_frame, [
            "Individual submission",
            "AgentClinic",
        ])

# ============================================================
# SLIDE 3 · Problem & Solution
# ============================================================
s3 = prs.slides[2]
for shape in s3.shapes:
    if not shape.has_text_frame:
        continue
    txt = shape.text_frame.text
    if "What real-world problem" in txt:
        set_multiline(shape.text_frame, [
            "AI agents ship to production without the testing discipline traditional software took decades to earn.",
            "Existing eval harnesses score outputs but cannot point at the trace event that drove the verdict — review becomes vibes, the same blind retry burns tokens release after release, debug is post-incident.",
            "Quality gates that work for traditional software do not cover AI agents.",
        ])
    elif "Brief summary of the solution" in txt:
        set_multiline(shape.text_frame, [
            "AgentClinic — a pre-production clinic for AI agents.",
            "Trace in, forensic report out. Every finding bound to specific trace_event_ids as evidence (no finding without an evidence span). Published natively into UiPath Test Cloud as Test Sets / Executions / per-pattern Case Logs.",
            "Judge (deterministic rule engine) and Coach (LLM, bounded) separated by a code-level contract. The LLM translates findings into remediation — it cannot revise the verdict. Coach, not surveillance.",
        ])

# ============================================================
# SLIDE 4 · Benefits & technologies used
# ============================================================
s4 = prs.slides[3]
# Table: 5 rows in fixed order — End-user / User department / Industries
# / UiPath products / Other tech. Use index, not label match (cells contain
# zero-width spaces that defeat startswith comparisons).
TABLE_VALUES_BY_ROW = [
    "AI agent platform engineers · ML platform teams · QA leads",
    "Engineering · Platform · SRE · Developer Experience",
    "Enterprise SaaS · AI infrastructure · any vertical deploying AI agents to production",
    "Test Cloud (Test Manager) · Coded Agent (Function) · Orchestrator Process · AgentHub LLM Gateway · AI Trust Layer · External Application",
    "Python 3.11+ · Pydantic · jsonschema · Claude Code (AgentHack Coding Agents bonus)",
]
for shape in s4.shapes:
    if shape.has_table:
        for i, row in enumerate(shape.table.rows):
            if i < len(TABLE_VALUES_BY_ROW):
                set_text(row.cells[1].text_frame, TABLE_VALUES_BY_ROW[i])
    elif shape.has_text_frame:
        txt = shape.text_frame.text
        if "What does this agent actually achieve" in txt:
            set_multiline(shape.text_frame, [
                "Quality shifts from late-stage checkpoint to a continuous, governed capability across every release.",
                "Every finding auditable via Test Cloud Case Log + evidence-bound override-result reason.",
                "No direct LLM provider API key — coach calls ride UiPath's AI Trust Layer for audit + PII redaction.",
                "Re-runnable on the same trace (idempotent per-pattern Test Cases).",
                "Production error matrix documented — 7 failure modes line-numbered in the repo.",
            ])
        elif shape.text_frame.text.strip() == "Details":
            set_text(shape.text_frame, "5 golden traces published live to Test Cloud project ACR2.")

# ============================================================
# SLIDE 5 · Solution architecture
# ============================================================
s5 = prs.slides[4]
for shape in s5.shapes:
    if not shape.has_text_frame:
        continue
    txt = shape.text_frame.text
    if "This slide is optional" in txt:
        set_multiline(shape.text_frame, [
            "Four roles, no overlap:",
            "JUDGE — deterministic Python rule engine, no LLM",
            "COACH — LLM via AgentHub LLM Gateway, bounded by validator (forbidden to judge)",
            "RECORDER — UiPath Test Cloud, Test Manager REST v2",
            "ORCHESTRATOR — UiPath Orchestrator Process + AI Trust Layer audit",
            "",
            "Pipeline: trace → normalize → detect (7 patterns) → score (L0–L3) → coach (optional) → publish",
            "",
            "Publish chain: Project (idempotent) → Test Case (per pattern, reusable) → Test Set (per publish) → Execution → Case Log + override-result reason + report.md attached",
        ])

# ============================================================
# SLIDE 6 · Miscellaneous → repurpose as Evidence Table
# ============================================================
s6 = prs.slides[5]
for shape in s6.shapes:
    if not shape.has_text_frame:
        continue
    txt = shape.text_frame.text
    if "Miscellaneous" == txt.strip():
        set_text(shape.text_frame, "Live evidence in UiPath Test Cloud (project ACR2)")
    elif "If you need extra slides" in txt:
        set_multiline(shape.text_frame, [
            "Five golden traces published with --coach uipath, covering score levels L0 through L3:",
            "",
            "01 hard_hat_loop — blind retry, state unchanged → Failed · 0.0 / L3 · Test Set ACR2:79",
            "02 clean_run — negative sample → Passed · 100 / L0 · Test Set ACR2:80",
            "03 redundant_and_full_read — 2 patterns fire → Failed×2 · 70 / L1 · Test Set ACR2:81",
            "04 unverified_claim_missing_tokens → Failed · 40 / L2 · Test Set ACR2:82",
            "05 interleaved_think_retry → Failed · 0.0 / L3 · Test Set ACR2:84",
            "",
            "All five live in Test Cloud as concrete artifacts — not paper claims.",
            "Each Test Case Log carries a multi-line evidence-bound reason that names the specific trace_event_ids driving the verdict, plus the full report.md as an attachment.",
        ])

# ============================================================
# SLIDE 7 · Closing
# ============================================================
s7 = prs.slides[6]
for shape in s7.shapes:
    if not shape.has_text_frame:
        continue
    txt = shape.text_frame.text
    if "Closing message" in txt:
        set_multiline(shape.text_frame, [
            "AgentClinic — coach, not surveillance.",
            "",
            "Drop a trace in. Get an audit trail.",
            "",
            "github.com/rainingsnow0914tw-ship-it/agentclinic",
            "UiPath AgentHack 2026 · Track 3 · Agentic Testing",
        ])

prs.save(DST)
print(f"Wrote {DST}")
print(f"Total slides: {len(prs.slides)}")
